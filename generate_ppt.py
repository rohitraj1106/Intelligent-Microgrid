from pptx import Presentation
from pptx.util import Inches, Pt
import os

def create_presentation(template_path, output_path, content_file):
    if not os.path.exists(template_path):
        print(f"Template not found: {template_path}")
        return

    prs = Presentation(template_path)
    
    # We will try to update the slides that already exist in the template
    # based on the content mapping. 
    # Slide index 0: Title
    # Slide index 1-15: Content
    
    # Simple Content Mapping (Slide Index -> Title, Body)
    content = [
        ("Intelligent Microgrid", "AI-Powered Energy Management\nProgress Update: Predictive Forecasting, Edge Data Layers, and P2P Trading"),
        ("Progress Overview", "- AI Backbone: Solar (2.84% MAPE) & Load (13.95% MAPE)\n- Edge Infrastructure: MQTT & SQLite telemetry\n- P2P Marketplace: Automated Order Book\n- Orchestration: FSM-based autonomous decisions"),
        ("Methodology", "1. NASA POWER API Acquisition\n2. Feature Engineering (Lags, Climate Data)\n3. XGBoost Hyper-tuned Modeling\n4. FSM Orchestration for Decisions"),
        ("Project Overview", "Objective: Decentralized energy management for N. India.\nScope:\n- Residential load/solar forecasting\n- Intelligent battery scheduling\n- P2P energy marketplace"),
        ("Data & Resources", "Datasets: Solar (175K rows), Load (3.28M rows)\nTech Stack: Python, XGBoost, MQTT, SQLite, FastAPI, NASA API"),
        ("Project Lifecycle", "Phase 1: Architecture\nPhase 2: Predictive Engine\nPhase 3: Integration (Current)\nPhase 4: Deployment & Dashboard (Next)"),
        ("Technical Updates — Architecture", "- Distributed Nodes (MQTT/SQLite)\n- Central Marketplace (FastAPI)\n- Strategist (FSM Orchestrator)\n- Resilient N-1 redundancy"),
        ("Results — Forecasting", "Solar MAPE: 2.84%\nLoad MAPE: 13.95%\nDaytime-focused solar accuracy; behavioral-driven load challenges."),
        ("Results — Edge Data Layer", "Demonstrated end-to-end data pipeline:\nSensor -> MQTT -> Ingestion -> SQLite.\nSub-second latency & 15-min summarization."),
        ("Results — P2P Marketplace", "Order Book with BUY/SELL automated matching.\nReal-time price discovery based on local surplus.\nProactive trading based on 24h forecasts."),
        ("Challenges & Solutions", "Problem: AC/Heater spikes. Solution: Weather-lags & feature drivers.\nProblem: Network instability. Solution: Local SQLite buffering."),
        ("Next Steps and Goals", "- 'Safe Window' battery logic\n- React-based UI Dashboard\n- LLM Agent integration"),
        ("Expected Results & Impact", "- 25% peak demand reduction\n- 15-20% user energy savings\n- Maximized local solar utilization"),
        ("Project Usecases & Scope", "Residential Societies, Rural Microgrids, Commercial Complexes.\nFuture: EV Charging integration."),
        ("Detailed Progress – Development", "Forecasting: Full XGBoost Sensitivity\nEdge: MQTT Nodes & Simulators\nMarketplace: API & Order Book\nOrchestrator: FSM Control"),
        ("THANK YOU", "GitHub: theabhinav0231/Intelligent-Microgrid\nQuestions?")
    ]

    for i, (title_text, body_text) in enumerate(content):
        if i < len(prs.slides):
            slide = prs.slides[i]
            # Try to find title and body placeholders
            for shape in slide.shapes:
                if shape.is_placeholder:
                    if shape.placeholder_format.type == 1 or shape.name.lower().startswith("title") or shape.name.lower().startswith("rectangle"): # Title
                        try:
                            shape.text = title_text
                        except: pass
                    elif shape.placeholder_format.type == 2 or shape.name.lower().startswith("textbox") or shape.name.lower().startswith("content"): # Body
                        try:
                            shape.text = body_text
                        except: pass

    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    base_dir = r"d:\Intelligent-Microgrid-main\Intelligent-Microgrid-main"
    template = os.path.join(base_dir, "f37c4ad2-0bf0-4d2a-be1c-28d33b79fc36.pptx")
    output = os.path.join(base_dir, "Intelligent_Microgrid_Progress_Report.pptx")
    content_file = os.path.join(base_dir, "presentation_content.md")
    
    create_presentation(template, output, content_file)
