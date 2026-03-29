import PyPDF2
from pptx import Presentation
import os

def extract_pdf_content(pdf_path):
    print(f"--- Extracting PDF: {pdf_path} ---")
    try:
        with open(pdf_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for i, page in enumerate(reader.pages):
                print(f"\nSlide {i+1}:")
                text = page.extract_text()
                if text:
                    print(text)
                else:
                    print("[No text found on this page]")
    except Exception as e:
        print(f"Error reading PDF: {e}")

def inspect_pptx_template(pptx_path):
    print(f"\n--- Inspecting PPTX Template: {pptx_path} ---")
    try:
        prs = Presentation(pptx_path)
        for i, slide in enumerate(prs.slides):
            print(f"\nSlide {i+1}:")
            # Print slide layout name
            print(f"Layout: {slide.slide_layout.name}")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    print(f"[{shape.name}] {shape.text}")
    except Exception as e:
        print(f"Error reading PPTX: {e}")

if __name__ == "__main__":
    # Get absolute paths to be safe
    base_dir = r"d:\Intelligent-Microgrid-main\Intelligent-Microgrid-main"
    pdf_file = os.path.join(base_dir, "minor project ppt.pdf")
    pptx_template = os.path.join(base_dir, "f37c4ad2-0bf0-4d2a-be1c-28d33b79fc36.pptx")
    
    extract_pdf_content(pdf_file)
    inspect_pptx_template(pptx_template)
